from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()

    title = spec.get("title")
    subtitle = spec.get("subtitle")
    if title or subtitle:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        if title:
            slide.shapes.title.text = str(title)
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = str(subtitle)

    for slide_spec in spec.get("slides", []):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(slide_spec.get("title", "Slide"))
        text_frame = slide.placeholders[1].text_frame
        bullets = slide_spec.get("bullets", [])
        if bullets:
            text_frame.text = str(bullets[0])
            for bullet in bullets[1:]:
                paragraph = text_frame.add_paragraph()
                paragraph.text = str(bullet)
        else:
            text_frame.text = str(slide_spec.get("body", ""))

    presentation.save(str(output_path))


if __name__ == "__main__":
    main()
